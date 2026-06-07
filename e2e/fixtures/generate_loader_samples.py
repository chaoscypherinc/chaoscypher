# Copyright (C) 2024-2026 Chaos Cypher, Inc.
# SPDX-License-Identifier: AGPL-3.0-only

"""Generate minimal valid sample files for the loader-matrix e2e test.

Run: ``python e2e/fixtures/generate_loader_samples.py``

One file per supported document loader. Files are intentionally tiny
(text content "Alice and Bob ..." so the matrix test can assert that
the parsed chunk content includes a known marker. Idempotent:
existing files are kept unless ``--force`` is passed.

The four binary fixtures (docx, xlsx, pptx, epub) replace the
sanitized-to-empty placeholders left by the public-repo prep
(``c9e47c918``) so that ``test_loader_matrix.py`` actually exercises
each loader against a parseable file. Tracked in TODO.md under
``Corrupt loader test fixtures``.
"""

from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path
from textwrap import dedent


SAMPLE_BODY = (
    "Alice and Bob are coworkers at Acme Corp. "
    "Alice manages the engineering team. "
    "Bob is a senior designer reporting to Alice."
)

OUTPUT_DIR = Path(__file__).parent / "sample_data"


# ---------------------------------------------------------------------------
# Plain-text formats
# ---------------------------------------------------------------------------


def write_markdown(path: Path) -> None:
    body = dedent(
        f"""\
        # Sample Markdown Document

        ## Section: People

        {SAMPLE_BODY}

        ## Section: Notes

        - Alice's role: Engineering Manager
        - Bob's role: Senior Designer
        """
    )
    path.write_text(body, encoding="utf-8")


def write_csv(path: Path) -> None:
    rows = [
        "name,role,team",
        "Alice,Engineering Manager,Engineering",
        "Bob,Senior Designer,Design",
    ]
    path.write_text("\n".join(rows) + "\n", encoding="utf-8")


def write_json(path: Path) -> None:
    data = {
        "people": [
            {"name": "Alice", "role": "Engineering Manager"},
            {"name": "Bob", "role": "Senior Designer"},
        ],
        "company": "Acme Corp",
        "summary": SAMPLE_BODY,
    }
    path.write_text(json.dumps(data, indent=2) + "\n", encoding="utf-8")


def write_jsonl(path: Path) -> None:
    lines = [
        {"name": "Alice", "role": "Engineering Manager"},
        {"name": "Bob", "role": "Senior Designer"},
    ]
    path.write_text(
        "\n".join(json.dumps(line) for line in lines) + "\n",
        encoding="utf-8",
    )


def write_html(path: Path) -> None:
    body = dedent(
        f"""\
        <!DOCTYPE html>
        <html lang="en">
          <head><meta charset="utf-8"><title>Sample</title></head>
          <body>
            <h1>Sample HTML Document</h1>
            <p>{SAMPLE_BODY}</p>
            <ul>
              <li>Alice — Engineering Manager</li>
              <li>Bob — Senior Designer</li>
            </ul>
          </body>
        </html>
        """
    )
    path.write_text(body, encoding="utf-8")


# ---------------------------------------------------------------------------
# Binary office formats (use the canonical Python libraries)
# ---------------------------------------------------------------------------


def write_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Sample DOCX Document", level=1)
    doc.add_heading("People", level=2)
    doc.add_paragraph(SAMPLE_BODY)
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Alice"
    table.rows[0].cells[1].text = "Engineering Manager"
    row = table.add_row()
    row.cells[0].text = "Bob"
    row.cells[1].text = "Senior Designer"
    doc.save(str(path))


def write_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "People"
    ws.append(["name", "role", "team"])
    ws.append(["Alice", "Engineering Manager", "Engineering"])
    ws.append(["Bob", "Senior Designer", "Design"])
    # Second sheet so multi-sheet handling is exercised.
    ws2 = wb.create_sheet("Summary")
    ws2.append(["summary"])
    ws2.append([SAMPLE_BODY])
    wb.save(str(path))


def write_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    # Title slide.
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Sample PPTX Document"
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "People at Acme Corp"
    # Body slide.
    body_slide = prs.slides.add_slide(prs.slide_layouts[1])
    body_slide.shapes.title.text = "People"
    body_slide.placeholders[1].text = (
        "Alice — Engineering Manager\nBob — Senior Designer\n\n" + SAMPLE_BODY
    )
    prs.save(str(path))


# ---------------------------------------------------------------------------
# EPUB — minimal valid bundle assembled by hand (no ebooklib dep).
# ---------------------------------------------------------------------------


_EPUB_CONTAINER = dedent(
    """\
    <?xml version="1.0" encoding="UTF-8"?>
    <container version="1.0"
        xmlns="urn:oasis:names:tc:opendocument:xmlns:container">
      <rootfiles>
        <rootfile full-path="OEBPS/content.opf"
                  media-type="application/oebps-package+xml"/>
      </rootfiles>
    </container>
    """
)

_EPUB_OPF = dedent(
    """\
    <?xml version="1.0" encoding="UTF-8"?>
    <package xmlns="http://www.idpf.org/2007/opf"
             unique-identifier="bookid" version="2.0">
      <metadata xmlns:dc="http://purl.org/dc/elements/1.1/"
                xmlns:opf="http://www.idpf.org/2007/opf">
        <dc:title>Sample EPUB Document</dc:title>
        <dc:language>en</dc:language>
        <dc:identifier id="bookid">urn:uuid:e2e-sample-epub</dc:identifier>
        <dc:creator>E2E</dc:creator>
      </metadata>
      <manifest>
        <item id="chapter1" href="chapter1.xhtml"
              media-type="application/xhtml+xml"/>
        <item id="ncx" href="toc.ncx"
              media-type="application/x-dtbncx+xml"/>
      </manifest>
      <spine toc="ncx">
        <itemref idref="chapter1"/>
      </spine>
    </package>
    """
)

_EPUB_NCX = dedent(
    """\
    <?xml version="1.0" encoding="UTF-8"?>
    <ncx xmlns="http://www.daisy.org/z3986/2005/ncx/" version="2005-1">
      <head><meta name="dtb:uid" content="urn:uuid:e2e-sample-epub"/></head>
      <docTitle><text>Sample EPUB Document</text></docTitle>
      <navMap>
        <navPoint id="n1" playOrder="1">
          <navLabel><text>Chapter 1</text></navLabel>
          <content src="chapter1.xhtml"/>
        </navPoint>
      </navMap>
    </ncx>
    """
)

_EPUB_CHAPTER = dedent(
    f"""\
    <?xml version="1.0" encoding="UTF-8"?>
    <!DOCTYPE html PUBLIC "-//W3C//DTD XHTML 1.1//EN"
        "http://www.w3.org/TR/xhtml11/DTD/xhtml11.dtd">
    <html xmlns="http://www.w3.org/1999/xhtml" xml:lang="en">
      <head><title>Chapter 1</title></head>
      <body>
        <h1>Chapter 1: People</h1>
        <p>{SAMPLE_BODY}</p>
        <ul>
          <li>Alice — Engineering Manager</li>
          <li>Bob — Senior Designer</li>
        </ul>
      </body>
    </html>
    """
)


def write_pdf(path: Path) -> None:
    """Hand-craft a minimal valid PDF with one extractable text line.

    ``pypdf.PdfWriter`` only knows how to add blank pages — building a
    page with text means writing a content stream with PDF operators
    (``BT … Tj … ET``), which means computing the xref table byte
    offsets. Easier to assemble the bytes directly than to depend on
    ``reportlab`` just for fixtures. Output is ~400 bytes.
    """
    # PDF objects assembled in order; we'll splice byte offsets into
    # the xref afterwards.
    text = (
        SAMPLE_BODY + " Roles: Alice is Engineering Manager. Bob is Senior Designer."
    )
    objects: list[bytes] = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 4 0 R "
            b">> >> /MediaBox [0 0 612 792] /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    # Object 5 is the content stream — text drawing.
    stream_body = (
        b"BT /F1 12 Tf 50 750 Td ("
        + text.encode("latin-1", errors="replace")
        + b") Tj ET"
    )
    stream_obj = (
        b"<< /Length " + str(len(stream_body)).encode() + b" >>\nstream\n"
        + stream_body
        + b"\nendstream"
    )
    objects.append(stream_obj)

    pdf = bytearray(b"%PDF-1.4\n%\xc4\xe5\xf2\xe5\xeb\xa7\xf3\xa0\xd0\xc4\xc6\n")
    offsets: list[int] = []
    for idx, obj in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf += f"{idx} 0 obj\n".encode()
        pdf += obj
        pdf += b"\nendobj\n"

    xref_offset = len(pdf)
    pdf += f"xref\n0 {len(objects) + 1}\n".encode()
    pdf += b"0000000000 65535 f \n"
    for off in offsets:
        pdf += f"{off:010d} 00000 n \n".encode()
    pdf += (
        b"trailer << /Size " + str(len(objects) + 1).encode()
        + b" /Root 1 0 R >>\n"
        + b"startxref\n"
        + str(xref_offset).encode()
        + b"\n%%EOF\n"
    )
    path.write_bytes(bytes(pdf))


def write_epub(path: Path) -> None:
    """Build a minimal valid EPUB v2 by hand.

    The ``mimetype`` entry MUST be the first file in the archive and
    MUST be stored uncompressed — that's how readers identify the
    bundle as EPUB without inspecting the rest.
    """
    with zipfile.ZipFile(path, "w") as zf:
        # mimetype: first entry, uncompressed
        zf.writestr(
            zipfile.ZipInfo("mimetype"),
            "application/epub+zip",
            compress_type=zipfile.ZIP_STORED,
        )
        zf.writestr("META-INF/container.xml", _EPUB_CONTAINER)
        zf.writestr("OEBPS/content.opf", _EPUB_OPF)
        zf.writestr("OEBPS/toc.ncx", _EPUB_NCX)
        zf.writestr("OEBPS/chapter1.xhtml", _EPUB_CHAPTER)


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------


SAMPLES: dict[str, callable] = {
    "sample.md": write_markdown,
    "sample.csv": write_csv,
    "sample.json": write_json,
    "sample.jsonl": write_jsonl,
    "sample.html": write_html,
    "sample.pdf": write_pdf,
    "sample.docx": write_docx,
    "sample.xlsx": write_xlsx,
    "sample.pptx": write_pptx,
    "sample.epub": write_epub,
}


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--force",
        action="store_true",
        help="Overwrite existing fixtures.",
    )
    args = parser.parse_args()

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    for filename, writer in SAMPLES.items():
        path = OUTPUT_DIR / filename
        if path.exists() and not args.force:
            print(f"  skip (exists): {filename}")
            continue
        writer(path)
        size = path.stat().st_size
        print(f"  wrote {filename} ({size} bytes)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
