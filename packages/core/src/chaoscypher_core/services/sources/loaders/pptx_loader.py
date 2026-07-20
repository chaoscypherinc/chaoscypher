# Copyright (C) 2024-2026 Chaos Cypher, Inc.
# SPDX-License-Identifier: AGPL-3.0-only

"""Microsoft PowerPoint (.pptx) loader.

Walks each slide's shapes, captures the title shape separately so it
can be surfaced in metadata, and concatenates body text frames in
shape order.
"""

from __future__ import annotations

import zipfile
from pathlib import Path
from typing import TYPE_CHECKING, Any

import structlog
from pptx import Presentation
from pptx.exc import PackageNotFoundError

from chaoscypher_core.exceptions import ValidationError
from chaoscypher_core.plugins import PluginMetadata


if TYPE_CHECKING:
    from chaoscypher_core.settings import EngineSettings


logger = structlog.get_logger(__name__)


class PPTXLoader:
    """Loader for Microsoft PowerPoint ``.pptx`` files."""

    @property
    def supported_extensions(self) -> list[str]:
        """File extensions this loader supports."""
        return [".pptx", ".PPTX"]

    @property
    def metadata(self) -> PluginMetadata:
        """Plugin metadata."""
        return PluginMetadata(
            plugin_id="pptx",
            name="PowerPoint PPTX Loader",
            description="Loads Microsoft PowerPoint .pptx files.",
            version="1.0.0",
            author="ChaosCypher",
            category="loader",
            builtin=True,
            origin="builtin",
            tags=["document", "office", "presentation"],
        )

    def __init__(self, settings: EngineSettings | None = None) -> None:
        """Initialize PPTX loader.

        Args:
            settings: Settings instance (not currently used).
        """
        self.settings = settings

    def load_document(self, filepath: str) -> list[dict[str, Any]]:
        """Load a PowerPoint presentation.

        Args:
            filepath: Path to a .pptx file.

        Returns:
            A single-item list with all slide text concatenated and
            metadata listing per-slide titles.

        Raises:
            ValidationError: If the file is not a valid PPTX (corrupt
                OPC package, missing zip central directory, etc.).
        """
        from chaoscypher_core.services.sources.loaders.base import check_loader_file_size

        check_loader_file_size(filepath, self.settings)

        path = Path(filepath)
        try:
            pres = Presentation(str(path))
        except (PackageNotFoundError, KeyError, zipfile.BadZipFile) as exc:
            msg = f"File '{path.name}' is not a valid PPTX: {exc}"
            raise ValidationError(msg, field="content") from exc

        parts: list[str] = []
        slide_summaries: list[dict[str, Any]] = []
        shapes_skipped_by_type: dict[str, int] = {}

        for idx, slide in enumerate(pres.slides, start=1):
            title = ""
            body_parts: list[str] = []
            title_shape = slide.shapes.title if slide.shapes.title is not None else None

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    type_name = (
                        str(shape.shape_type).split(" ")[0]
                        if shape.shape_type is not None
                        else "UNKNOWN"
                    )
                    shapes_skipped_by_type[type_name] = shapes_skipped_by_type.get(type_name, 0) + 1
                    continue
                is_title_shape = title_shape is not None and shape == title_shape
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if not text:
                        continue
                    if is_title_shape and not title:
                        title = text
                    else:
                        body_parts.append(text)

            slide_text_parts: list[str] = []
            if title:
                slide_text_parts.append(f"\n[Slide {idx}: {title}]")
            else:
                slide_text_parts.append(f"\n[Slide {idx}]")
            slide_text_parts.extend(body_parts)
            parts.append("\n".join(slide_text_parts))
            slide_summaries.append({"index": idx, "title": title})

        content = "\n\n".join(parts)

        logger.info(
            "pptx_loaded",
            filepath=str(path),
            slide_count=len(pres.slides),
            shapes_skipped_by_type=shapes_skipped_by_type,
            character_count=len(content),
        )

        return [
            {
                "content": content,
                "metadata": {
                    "source": str(path),
                    "extraction_method": "pptx",
                    "slides": slide_summaries,
                    "slide_count": len(pres.slides),
                    "content_type": (
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    ),
                    "loader_pptx_shapes_skipped": shapes_skipped_by_type,
                },
            }
        ]

    def supports_ocr(self) -> bool:
        """PPTX files don't need OCR."""
        return False


__all__ = ["PPTXLoader"]
