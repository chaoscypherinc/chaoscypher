# Copyright (C) 2024-2026 Chaos Cypher, Inc.
# SPDX-License-Identifier: AGPL-3.0-only

"""PPTXLoader extracts slide titles and body text frames.

Workstream 7 (2026-05-07): standalone .pptx uploads had no loader.
"""

from __future__ import annotations

import io
import struct
import zlib
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from chaoscypher_core.exceptions import ValidationError
from chaoscypher_core.services.sources.loaders.pptx_loader import PPTXLoader


FIXTURE = Path(__file__).parent.parent.parent.parent.parent / "fixtures" / "loaders" / "sample.pptx"


# ---------------------------------------------------------------------------
# Fixture helpers
# ---------------------------------------------------------------------------


def _minimal_png_bytes() -> bytes:
    """Return a minimal valid 1x1 white PNG as bytes."""

    def chunk(name: bytes, data: bytes) -> bytes:
        c = struct.pack(">I", len(data)) + name + data
        crc = zlib.crc32(c[4:]) & 0xFFFFFFFF
        return c + struct.pack(">I", crc)

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    idat_comp = zlib.compress(b"\x00\xff\xff\xff")
    idat = chunk(b"IDAT", idat_comp)
    iend = chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


def _build_pptx_with_image_and_text(tmp_path: Path) -> Path:
    """Return a .pptx file with one PICTURE shape and one text shape.

    The picture shape has ``has_text_frame=False`` so it will be counted
    in ``loader_pptx_shapes_skipped`` under the key ``"PICTURE"``.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    img_io = io.BytesIO(_minimal_png_bytes())
    slide.shapes.add_picture(img_io, Inches(1), Inches(1), Inches(1), Inches(1))

    txbox = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(3), Inches(1))
    txbox.text_frame.text = "Hello world"

    out = tmp_path / "with_image.pptx"
    prs.save(str(out))
    return out


def _build_pptx_with_text_only(tmp_path: Path) -> Path:
    """Return a .pptx file that contains only text-frame shapes (nothing skipped)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    txbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txbox.text_frame.text = "Only text here"

    out = tmp_path / "text_only.pptx"
    prs.save(str(out))
    return out


def test_pptx_loader_extracts_slide_titles_and_body() -> None:
    assert FIXTURE.exists(), f"fixture missing: {FIXTURE}"
    loader = PPTXLoader()
    docs = loader.load_document(str(FIXTURE))
    text = docs[0]["content"]
    assert "Slide 1 Title" in text
    assert "First slide bullet point" in text
    assert "Second Slide" in text
    assert "Another bullet" in text


def test_pptx_loader_records_slide_summaries() -> None:
    loader = PPTXLoader()
    docs = loader.load_document(str(FIXTURE))
    md = docs[0]["metadata"]
    assert md["slide_count"] == 2
    assert md["slides"][0]["title"] == "Slide 1 Title"
    assert md["slides"][1]["title"] == "Second Slide"
    assert md["extraction_method"] == "pptx"


def test_pptx_loader_supported_extensions() -> None:
    loader = PPTXLoader()
    assert ".pptx" in loader.supported_extensions
    assert loader.metadata.plugin_id == "pptx"


def test_pptx_loader_wraps_invalid_file_in_validation_error(tmp_path: Path) -> None:
    """A file with .pptx extension but wrong magic bytes raises ValidationError.

    Without the wrapper, python-pptx surfaces ``PackageNotFoundError``
    or a ``zipfile.BadZipFile`` — both bypass the API exception envelope.
    """
    bogus = tmp_path / "fake.pptx"
    bogus.write_bytes(b"not a real pptx")

    loader = PPTXLoader()
    with pytest.raises(ValidationError) as exc_info:
        loader.load_document(str(bogus))

    err = exc_info.value
    assert err.code == "VALIDATION_ERROR"
    assert err.field == "content"
    assert "fake.pptx" in err.message
    assert "not a valid" in err.message
    assert exc_info.value.__cause__ is not None


def test_pptx_loader_emits_per_type_dict(tmp_path: Path) -> None:
    """PPTXLoader records skipped shapes by python-pptx shape_type name (Phase 7
    audit-remediation 2026-05-09).
    """
    pptx_path = _build_pptx_with_image_and_text(tmp_path)

    loader = PPTXLoader()
    docs = loader.load_document(str(pptx_path))

    skipped = docs[0]["metadata"].get("loader_pptx_shapes_skipped")
    assert isinstance(skipped, dict)
    # python-pptx reports add_picture shapes as "PICTURE"
    assert "PICTURE" in skipped
    assert skipped["PICTURE"] == 1
    assert sum(skipped.values()) >= 1


def test_pptx_loader_emits_empty_dict_when_no_shapes_skipped(tmp_path: Path) -> None:
    """When no shapes are skipped, metadata holds {} (empty dict), not None."""
    pptx_path = _build_pptx_with_text_only(tmp_path)

    loader = PPTXLoader()
    docs = loader.load_document(str(pptx_path))

    skipped = docs[0]["metadata"].get("loader_pptx_shapes_skipped")
    assert skipped == {}
